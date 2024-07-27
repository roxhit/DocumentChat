import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

load_dotenv()
google_api_key = os.getenv("GOOGLE_API_KEY") 

if not google_api_key:
    raise ValueError("GOOGLE_API_KEY environment variable not found")

genai.configure(api_key=google_api_key)

def get_pdf_text(pdf):
    text = ""
    pdf_reader = PdfReader(pdf)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_docx_text(docx):
    doc = Document(docx)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def get_pptx_text(pptx):
    prs = Presentation(pptx)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def get_xlsx_text(xlsx):
    wb = load_workbook(xlsx)
    text = ""
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
    return text

def get_text_from_file(file):
    if file.name.endswith('.pdf'):
        return get_pdf_text(file)
    elif file.name.endswith('.docx'):
        return get_docx_text(file)
    elif file.name.endswith('.pptx'):
        return get_pptx_text(file)
    elif file.name.endswith('.xlsx'):
        return get_xlsx_text(file)
    else:
        return ""

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    google_api_key = os.getenv("GOOGLE_API_KEY")
    if not google_api_key:
        raise ValueError("GOOGLE_API_KEY environment variable not found")
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=google_api_key)
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    google_api_key = os.getenv("GOOGLE_API_KEY")
    if not google_api_key:
        raise ValueError("GOOGLE_API_KEY environment variable not found")
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=google_api_key)
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    return response["output_text"]

def main():
    st.set_page_config("Chat with Documents")
    st.header("Chat with your Documents")
    if "messages" not in st.session_state:
        st.session_state.messages = []

    with st.sidebar:
        st.title("Menu:")
        uploaded_files = st.file_uploader("Upload your Files and Click on the Submit & Process Button", accept_multiple_files=True)
        if st.button("Submit & Process", key="process"):
            with st.spinner("Processing..."):
                raw_text = ""
                for file in uploaded_files:
                    raw_text += get_text_from_file(file)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("Done")

    with st.form("input_form"):
        user_question = st.text_input("Ask a Question from the Uploaded Files")
        if st.form_submit_button("Send"):
            if user_question:
                st.session_state.messages.append({"role": "user", "content": user_question})
                with st.spinner("Processing..."):
                    response = user_input(user_question)
                    st.session_state.messages.append({"role": "bot", "content": response})

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

if __name__ == "__main__":
    main()